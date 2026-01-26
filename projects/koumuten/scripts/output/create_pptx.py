from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# プレゼンテーション作成
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """タイトルスライドを追加"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets):
    """コンテンツスライドを追加"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True

    # 箇条書き
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.space_after = Pt(12)

    return slide

def add_table_slide(prs, title, headers, rows):
    """テーブルスライドを追加"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True

    # テーブル
    cols = len(headers)
    rows_count = len(rows) + 1
    table = slide.shapes.add_table(rows_count, cols, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.5 * rows_count)).table

    # ヘッダー
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(18)

    # データ
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(16)

    return slide

# スライド1: 表紙
add_title_slide(prs, "工務店専門LP制作サービス", "Web集客で問い合わせ数を2倍に")

# スライド2: こんなお悩みありませんか？
add_content_slide(prs, "こんなお悩みありませんか？", [
    "チラシ・ポータルサイト依存から抜け出せない",
    "ホームページはあるが問い合わせが来ない",
    "見学会の集客に苦戦している",
    "Web広告を出しても成果が出ない",
    "SNSをやっているが問い合わせに繋がらない"
])

# スライド3: 工務店業界の現状
add_table_slide(prs, "工務店業界の現状",
    ["指標", "数値"],
    [
        ["新設住宅着工戸数", "前年比 -4.6%"],
        ["持家着工戸数", "前年比 -11.4%"],
        ["工務店数", "約40万社"],
        ["Web集客対応", "多くが未対応"]
    ])

# スライド4: なぜLPが必要なのか
add_table_slide(prs, "HPとLPの違い",
    ["項目", "ホームページ", "LP"],
    [
        ["目的", "会社紹介", "1つの行動を促す"],
        ["情報量", "多い", "絞られている"],
        ["離脱率", "高い", "低い"],
        ["CV率", "1-2%", "3-10%"]
    ])

# スライド5: LPの効果（実績）
add_content_slide(prs, "LPの効果（実績）", [
    "【完成見学会LP】来場3組 → 9組（3倍）",
    "【リフォームLP】月3件 → 月15件（5倍）",
    "【採用LP】応募0件 → 月5件",
    "",
    "→ LPは「問い合わせ獲得マシン」"
])

# スライド6: 私たちの強み
add_content_slide(prs, "私たちの強み", [
    "【工務店専門】業界用語・慣習を理解。工務店特有の訴求ポイントを熟知",
    "【成果にコミット】問い合わせ保証プランあり。公開後の改善提案",
    "【スピード対応】最短2週間で公開。見学会に間に合う短納期対応"
])

# スライド7: サービス内容
add_content_slide(prs, "サービス内容", [
    "競合調査・市場分析",
    "構成案作成（ワイヤーフレーム）",
    "デザイン制作",
    "コーディング・実装（スマホ対応）",
    "フォーム設置",
    "公開後1ヶ月サポート"
])

# スライド8: 制作の流れ
add_content_slide(prs, "制作の流れ", [
    "1週目: ヒアリング・競合調査",
    "2週目: 構成案・デザイン案作成",
    "3週目: デザイン確定・コーディング",
    "4週目: テスト・公開",
    "公開後: 効果測定・改善提案"
])

# スライド9: 料金プラン
add_table_slide(prs, "料金プラン",
    ["プラン", "内容", "価格"],
    [
        ["スタンダード", "LP1ページ", "○○万円〜"],
        ["プレミアム", "LP + 広告運用", "○○万円〜/月"],
        ["カスタム", "ご要望に応じて", "要相談"]
    ])

# スライド10: よくある質問
add_content_slide(prs, "よくある質問", [
    "Q: 自社でも更新できますか？ → はい、更新マニュアルをお渡しします",
    "Q: 広告運用もお願いできますか？ → はい、Google/Meta広告に対応",
    "Q: 既存HPとの連携は？ → サブドメインまたはサブディレクトリで対応可能"
])

# スライド11: 次のステップ
add_content_slide(prs, "次のステップ", [
    "【無料相談のご案内】",
    "",
    "1. 現状ヒアリング（15分）",
    "2. 改善ポイントのご提示",
    "3. ご提案書の作成",
    "",
    "→ まずはお気軽にご相談ください"
])

# スライド12: お問い合わせ
add_title_slide(prs, "お問い合わせ", "Email: xxxx@example.com\nTEL: 000-0000-0000")

# 保存
output_path = "projects/koumuten/output/koumuten_lp_sales_deck.pptx"
prs.save(output_path)
print(f"パワーポイントファイルを作成しました: {output_path}")
